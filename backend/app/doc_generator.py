from io import BytesIO
from typing import Tuple
from docx import Document
from pptx import Presentation

from . import models


def _build_docx(project: models.Project) -> BytesIO:
    doc = Document()
    doc.add_heading(project.title, level=1)
    sections = sorted(project.sections, key=lambda s: s.order_index)
    for sec in sections:
        doc.add_heading(sec.title, level=2)
        text = sec.content or ""
        for block in text.split("\n\n"):
            block = block.strip()
            if block:
                doc.add_paragraph(block)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _build_pptx(project: models.Project) -> BytesIO:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = project.title
    if slide.placeholders and len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""

    bullet_layout = prs.slide_layouts[1]
    sections = sorted(project.sections, key=lambda s: s.order_index)
    for sec in sections:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = sec.title
        body = slide.placeholders[1].text_frame
        body.clear()
        text = (sec.content or "").strip()
        if not text:
            continue
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        first = True
        for line in lines:
            if first:
                body.text = line
                first = False
            else:
                p = body.add_paragraph()
                p.text = line
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_document_file(project: models.Project) -> Tuple[BytesIO, str, str]:
    ptype = (project.project_type or "").lower()
    safe_title = project.title.replace(" ", "_")
    if ptype == "docx":
        buf = _build_docx(project)
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        filename = f"{safe_title}.docx"
        return buf, mime, filename
    if ptype == "pptx":
        buf = _build_pptx(project)
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        filename = f"{safe_title}.pptx"
        return buf, mime, filename
    raise ValueError(f"Unsupported project_type: {project.project_type}")
